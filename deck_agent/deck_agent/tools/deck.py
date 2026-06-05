"""
Deck tools: validate_spec and render_deck.

validate_spec wraps the layered validator, supplying the data catalog so the
renderability layer can check shapes/capacities.

render_deck is BELT-AND-SUSPENDERS: it re-validates internally and refuses to
render an invalid spec, regardless of what the model did. "Render only valid
specs" is thus a property of the SYSTEM, not a behavior we trust the model to
follow. The model is TOLD to validate first (efficiency: catch errors before
attempting render); the renderer ENFORCES it (safety).

The actual python-pptx rendering is STUBBED — build it against your frozen spec
schema. The render-time precondition assertions are sketched: even past
validation, assert each element fits its region as you place it, so any
validator/renderer DRIFT surfaces loudly rather than producing a broken file.
"""

from __future__ import annotations

from datetime import datetime
from pathlib import Path

from ..validation import validate_spec as _validate
from .data import fetch_data_impl, get_data_catalog


def validate_spec_impl(spec: dict) -> dict:
    catalog = get_data_catalog()
    result = _validate(spec, catalog)
    return result.to_tool_result()


def render_deck_impl(spec: dict, output_name: str, output_dir: str | None = None) -> dict:
    # Belt-and-suspenders: never render an invalid spec, no matter what.
    catalog = get_data_catalog()
    result = _validate(spec, catalog)
    if not result.ok:
        # Surface as an error the model can act on, same shape as validate_spec.
        return {
            "rendered": False,
            "reason": "Spec failed validation at render time.",
            **result.to_tool_result(),
        }

    out_dir = Path(output_dir or "output")
    out_dir.mkdir(parents=True, exist_ok=True)
    stamp = datetime.now().strftime("%Y-%m-%d")
    out_path = out_dir / f"{output_name}_{stamp}.pptx"

    _render(spec, out_path)  # STUB — see below

    return {"rendered": True, "output_path": str(out_path)}


# ---------------------------------------------------------------------------
# STUB renderer. Replace with your audited python-pptx implementation.
# ---------------------------------------------------------------------------

def _render(spec: dict, out_path: Path) -> None:
    """Build the .pptx from a validated spec.

    TODO: implement against your frozen schema. Sketch of the intended shape:

        from pptx import Presentation
        prs = Presentation("templates/base_template.pptx")
        for slide_spec in spec["slides"]:
            layout = _resolve_layout(prs, slide_spec["layout"])
            slide = prs.slides.add_slide(layout)
            for element in slide_spec["elements"]:
                _place_element(slide, element)   # <-- assert preconditions here
        prs.save(out_path)

    Render-time precondition assertions (drift backstop): inside _place_element,
    before placing a table, assert its row count fits the region; before placing
    a chart image, assert its dimensions fit chart_region; raise a clear error if
    reality contradicts what validation predicted. That error, fed back through
    the loop, both fixes the run AND signals that validator/renderer have drifted.
    """
    # Minimal stub so the smoke test produces a file artifact.
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        for slide_spec in spec["slides"]:
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank-ish
            title = _first_title(slide_spec)
            if slide.shapes.title is not None:
                slide.shapes.title.text = title
            box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4))
            tf = box.text_frame
            tf.text = f"[STUB] layout={slide_spec.get('layout')} " \
                      f"elements={len(slide_spec.get('elements', []))}"
            tf.paragraphs[0].font.size = Pt(14)
        prs.save(out_path)
    except ImportError:
        # python-pptx not installed — write a placeholder so the loop completes.
        out_path.write_text("python-pptx not installed; stub output.\n")


def _first_title(slide_spec: dict) -> str:
    for el in slide_spec.get("elements", []):
        if el.get("type") == "title":
            return el.get("text", "Untitled")
    return "Untitled"
